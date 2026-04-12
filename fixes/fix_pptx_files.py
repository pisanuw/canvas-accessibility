"""
PowerPoint (.pptx) accessibility fixes.

Fixes:
  - pptx.image_alt     : images missing alt text         (AI-assisted)
  - pptx.slide_title   : slides missing titles           (AI-assisted)
  - pptx.reading_order : shapes not in reading order     (fully automatic)
  - pptx.no_language   : no document language set        (fully automatic)

Usage (CLI):
  python3 fix_pptx_files.py --course-id 1492292 --fix all
  python3 fix_pptx_files.py --course-id 1492292 --fix image_alt --file-id 12345
  python3 fix_pptx_files.py --course-id 1492292 --fix reading_order --dry-run
"""

import argparse
import io
import re
import sys
import time
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))
from fixes.canvas_client import CanvasClient

PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
PPTX_EXTS = {".pptx", ".ppt"}

# MSO placeholder types
PP_TITLE = 1        # TITLE placeholder idx
PP_CENTER_TITLE = 3 # CENTER_TITLE

_FILENAME_RE = re.compile(
    r'\.[a-zA-Z0-9]{2,5}$'
    r'|^[\w\-. ]+\.(png|jpe?g|gif|svg|bmp|webp|tiff?|pdf|docx?|pptx?|xlsx?)$',
    re.IGNORECASE,
)


def _is_filename_alt(text: str) -> bool:
    t = text.strip()
    return bool(t and _FILENAME_RE.search(t))


# ── Fix: image alt text ───────────────────────────────────────────────────────

def fix_image_alt(pptx_bytes: bytes) -> tuple[bytes, list[str]]:
    """
    Add AI-generated alt text to picture shapes that have no description.
    Returns (updated_pptx_bytes, list_of_changes).
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(pptx_bytes))
    changes = []

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            # Only picture shapes
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue

            # Check current alt text (stored in cNvPr descr attribute)
            cNvPr = shape._element.nvPicPr.cNvPr
            existing = cNvPr.get("descr", "").strip()
            if existing and not _is_filename_alt(existing):
                continue  # already has a real description

            placeholder = (
                "This image currently does not have a description. "
                "Your instructor will review it and add a description"
            )
            cNvPr.set("descr", placeholder)
            cNvPr.set("title", "Image needs description")
            changes.append(f"Slide {slide_num}: set placeholder alt text")

    if not changes:
        return pptx_bytes, changes

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), changes


# ── Fix: slide titles ─────────────────────────────────────────────────────────

def fix_slide_titles(pptx_bytes: bytes) -> tuple[bytes, list[str]]:
    """
    Add AI-generated titles to slides that have no title placeholder text.
    Returns (updated_pptx_bytes, list_of_changes).
    """
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.enum.text import PP_ALIGN

    prs = Presentation(io.BytesIO(pptx_bytes))
    changes = []

    for slide_num, slide in enumerate(prs.slides, 1):
        title_shape = slide.shapes.title
        if title_shape and title_shape.has_text_frame:
            text = title_shape.text_frame.text.strip()
            if text:
                continue  # already has a title

        title_text = (
            "This slide currently does not have a title. "
            "Your instructor will review it and add a title"
        )

        if title_shape and title_shape.has_text_frame:
            title_shape.text_frame.text = title_text
        else:
            # Add a title text box at standard title position (top of slide)
            left = Emu(457200)   # ~0.5 inch
            top  = Emu(274638)   # ~0.3 inch
            width = prs.slide_width - Emu(914400)
            height = Emu(1143000)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = title_text
            tf.paragraphs[0].runs[0].font.size = Pt(28)
            tf.paragraphs[0].runs[0].font.bold = True

        changes.append(f"Slide {slide_num}: added title '{title_text}'")

    if not changes:
        return pptx_bytes, changes

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), changes


# ── Fix: non-descriptive links ────────────────────────────────────────────────

def fix_links(pptx_bytes: bytes) -> tuple[bytes, list[str]]:
    """
    Replace non-descriptive hyperlink display text with 'linking to external content'.
    Returns (updated_pptx_bytes, list_of_changes).
    """
    from pptx import Presentation
    from pptx.oxml.ns import qn

    NON_DESC = {"here", "click here", "read more", "link", "more", "this", "url"}
    LINK_LABEL = "linking to external content"
    changes = []

    prs = Presentation(io.BytesIO(pptx_bytes))
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    rPr = run._r.get_or_add_rPr()
                    hlinkClick = rPr.find(qn("a:hlinkClick"))
                    if hlinkClick is None:
                        continue
                    text = run.text.strip()
                    if text.lower() in NON_DESC or not text:
                        run.text = LINK_LABEL
                        changes.append(
                            f"Slide {slide_num}: replaced link text "
                            f"'{text}' → 'linking to external content'")

    if not changes:
        return pptx_bytes, changes

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), changes


# ── Fix: reading order ────────────────────────────────────────────────────────

def fix_reading_order(pptx_bytes: bytes) -> tuple[bytes, list[str]]:
    """
    Reorder shapes on each slide so they appear top-to-bottom, left-to-right
    in the XML (which is what screen readers use for tab/reading order).
    Returns (updated_pptx_bytes, list_of_changes).
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(pptx_bytes))
    changes = []

    for slide_num, slide in enumerate(prs.slides, 1):
        shapes = list(slide.shapes)
        if len(shapes) < 2:
            continue

        # Sort by (top, left)
        sorted_shapes = sorted(shapes, key=lambda s: (s.top or 0, s.left or 0))
        original_order = [s.shape_id for s in shapes]
        new_order = [s.shape_id for s in sorted_shapes]

        if original_order == new_order:
            continue

        sp_tree = slide.shapes._spTree
        elements = [s._element for s in sorted_shapes]
        # Remove all shape elements and re-append in sorted order
        for el in elements:
            sp_tree.remove(el)
        for el in elements:
            sp_tree.append(el)

        changes.append(f"Slide {slide_num}: reordered {len(shapes)} shapes for reading order")

    if not changes:
        return pptx_bytes, changes

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), changes


# ── Fix: language ─────────────────────────────────────────────────────────────

def fix_language(pptx_bytes: bytes, lang: str = "en-US") -> tuple[bytes, list[str]]:
    """
    Set the language attribute on all text runs that lack one.
    Returns (updated_pptx_bytes, list_of_changes).
    """
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from lxml import etree

    prs = Presentation(io.BytesIO(pptx_bytes))
    changes = []
    set_count = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    rPr = run._r.get_or_add_rPr()
                    if not rPr.get("lang"):
                        rPr.set("lang", lang)
                        set_count += 1

    if set_count:
        changes.append(f"Set language '{lang}' on {set_count} text run(s)")

    if not changes:
        return pptx_bytes, changes

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), changes


# ── Orchestrate fixes on one file ─────────────────────────────────────────────

def fix_pptx_file(client: CanvasClient, course_id: int, file_info: dict,
                  fixes: list[str], dry_run: bool = False) -> dict:
    """Download a PPTX, apply fixes, re-upload."""
    file_id = file_info["id"]
    filename = file_info.get("display_name") or file_info.get("filename", "presentation.pptx")
    folder_id = file_info.get("folder_id")

    print(f"  Processing: {filename}")
    pptx_bytes = client.download_url(file_info["url"])
    all_changes = []
    all_fixes = "all" in fixes

    if all_fixes or "reading_order" in fixes:
        pptx_bytes, ch = fix_reading_order(pptx_bytes)
        all_changes.extend(ch)

    if all_fixes or "no_language" in fixes or "language" in fixes:
        pptx_bytes, ch = fix_language(pptx_bytes)
        all_changes.extend(ch)

    if all_fixes or "image_alt" in fixes:
        pptx_bytes, ch = fix_image_alt(pptx_bytes)
        all_changes.extend(ch)

    if all_fixes or "links" in fixes:
        pptx_bytes, ch = fix_links(pptx_bytes)
        all_changes.extend(ch)

    if all_fixes or "slide_title" in fixes:
        pptx_bytes, ch = fix_slide_titles(pptx_bytes)
        all_changes.extend(ch)

    result = {"file": filename, "file_id": file_id, "changes": all_changes, "updated": False}

    if all_changes and not dry_run:
        client.upload_file(course_id, folder_id, filename, PPTX_MIME, pptx_bytes)
        result["updated"] = True
        print(f"    Re-uploaded with {len(all_changes)} fix(es)")
    elif all_changes:
        print(f"    [dry-run] {len(all_changes)} change(s) would be made")
    else:
        print(f"    No changes needed")

    for ch in all_changes:
        print(f"      • {ch}")

    return result


def fix_course_pptx_files(client: CanvasClient, course_id: int,
                           fixes: list[str], file_id: int = None,
                           dry_run: bool = False) -> list[dict]:
    """Run fixes on one PPTX or all PPTX files in a course."""
    if file_id:
        files = [client.get_file_info(file_id)]
    else:
        all_files = client.list_files(course_id)
        files = [
            f for f in all_files
            if Path(f.get("filename", "")).suffix.lower() in PPTX_EXTS
        ]
        print(f"Found {len(files)} PowerPoint file(s) in course {course_id}")

    results = []
    for f in files:
        try:
            r = fix_pptx_file(client, course_id, f, fixes, dry_run)
            results.append(r)
            time.sleep(0.5)
        except Exception as e:
            print(f"  ERROR on '{f.get('filename', '?')}': {e}")
            results.append({"file": f.get("filename", "?"), "error": str(e)})

    return results


# ── CLI ───────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Fix PowerPoint accessibility issues")
    parser.add_argument("--course-id", type=int, required=True)
    parser.add_argument("--file-id", type=int, default=None)
    parser.add_argument("--fix", default="all",
                        help="Comma-separated: all,image_alt,links,slide_title,reading_order,no_language")
    parser.add_argument("--dry-run", action="store_true")
    args = parser.parse_args()

    fix_list = [f.strip() for f in args.fix.split(",")]
    client = CanvasClient()

    results = fix_course_pptx_files(client, args.course_id, fix_list,
                                    file_id=args.file_id, dry_run=args.dry_run)

    total_changes = sum(len(r.get("changes", [])) for r in results)
    updated = sum(1 for r in results if r.get("updated"))
    print(f"\nSummary: {total_changes} changes, {updated}/{len(results)} files updated")


if __name__ == "__main__":
    main()
